# !/usr/bin/python3
# -*- coding: utf-8 -*-
# @Time : 2021/3/24 11:00
# @Author : Mat
# @Email : mat_wu@163.com
# @File : erweima_detection.py
# @Software: PyCharm




# readfile = r"D:\20210202-康师傅×linefriends新品上市活动.pptx"
# img_folder = r"D:\\op"
# utils.save_pptx_as_png(img_folder,readfile,overwrite_folder=True)


import os
import shutil
import pyzbar.pyzbar as pyzbar
import cv2
from pptx import Presentation
from PyPDF2 import PdfFileReader, PdfFileWriter
from pdf2image import convert_from_path
from pptx_tools import utils
import argparse
from tqdm import tqdm

# 读取源文件文件名
def readSources(path):
    fileSource_list = os.listdir(path)
    return fileSource_list

# 将源文件转换为图片以方便检测二维码
def sourceToimg(n,filepath,outputdirectory,num_workers):
    if not os.path.exists(outputdirectory):
        os.mkdir(outputdirectory)
    # if 'pdf' in filepath:
    #     output = outputdirectory.split('/')[-1]
    #     convert_from_path(filepath, 300, output, fmt="JPEG", output_file="ok",
    #                       thread_count=num_workers)
    # elif 'pptx' in filepath:
    #     utils.save_pptx_as_png(n+'/'+outputdirectory.split('.')[-1],filepath,overwrite_folder=True)
    #     image_path = os.listdir(outputdirectory)
    #     for i in tqdm(image_path):
    #         newname = i.replace(i, 'ppt_erwei-%d' % int(i.split('片')[-1].split('.PNG')[0]) + '.jpg')
    #         os.rename(os.path.join(outputdirectory, i), os.path.join(outputdirectory, newname))

# 二维码检测
def qrcodeDetection(outputdirectory):
    list_erweimaindex = []
    img_path = os.listdir(outputdirectory)
    for i in tqdm(img_path):
        cvm = cv2.imread(outputdirectory + '\\' + i)
        gray = cv2.cvtColor(cvm, cv2.COLOR_BGR2GRAY)
        barcodes = pyzbar.decode(gray)
        for barcode in barcodes:
            barcodeData = barcode.data.decode("utf-8")
            print(barcodeData)
            if barcodeData is not None:
                list_erweimaindex.append(int(i.split('-')[-1].split('.jpg')[0]) - 1)
    list_erweimaindex.sort()
    shutil.rmtree(outputdirectory)
    print(list_erweimaindex)
    return list_erweimaindex

# 删除某一页ppt
def del_slide(prs, index):
    slides = list(prs.slides._sldIdLst)
    prs.slides._sldIdLst.remove(slides[index])

def generateSources(sourcefile,list_erweimaindex):
    if 'pdf' in sourcefile:
        pdfReader = PdfFileReader(open(sourcefile, 'rb'))
        pdfFileWriter = PdfFileWriter()
        numPages = pdfReader.getNumPages()
        pagelist = list_erweimaindex  # 注意第一页的index为0.
        for index in tqdm(range(0, numPages)):
            if index not in pagelist:
                pageObj = pdfReader.getPage(index)
                pdfFileWriter.addPage(pageObj)
        pdfFileWriter.write(open(sourcefile.split('.')[0]+'(删除二维码后文件)'+'.pdf', 'wb'))
    elif 'pptx' in sourcefile:
        ppt = Presentation(sourcefile)
        # 获取所有页
        slides = ppt.slides
        number_pages = len(slides)
        print("删除前ppt一共", number_pages, "页面")
        count = 0
        for index in tqdm(list_erweimaindex):
            # index1 = index-count
            # print(index1)
            del_slide(ppt, index - count)
            count += 1

        # 再次获取所有页
        slides = ppt.slides
        number_pages = len(slides)
        print("删除后ppt一共", number_pages, "页面")

        ppt.save(sourcefile.split('.')[0]+'(删除二维码后文件)'+'.pptx')
        print('生成完毕')

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description=' Project for qrcode detection and remove related pages with PPT and PDF')
    parser.add_argument('--num_workers', type=int, default=6,
                        help='thread for changing the pdf pages into pdf')  # 设置pdf转图片页面的线程数
    parser.add_argument('--pages_convert_image_output', type=str, default='./output',
                        help='pages convert image output')  # 图片输出文件夹
    parser.add_argument('--PPTAndPDF_directory', type=str, default='D:\\uyu',
                        help='sources of ppt and pdf with qrcode')  # ppt和pdf源文件
    args = parser.parse_args()  # 解析添加参数
    #os.mkdir(args.pages_convert_image_output)

    # 读取源文件
    fileSource_list = readSources(args.PPTAndPDF_directory)
    for file in tqdm(fileSource_list):
        file_name = args.PPTAndPDF_directory + '//' + file
        print(file_name)
        sourceToimg(args.PPTAndPDF_directory,file_name, args.pages_convert_image_output, args.num_workers)
        # list_erweimaindex = qrcodeDetection(args.pages_convert_image_output)
        # generateSources(file_name, list_erweimaindex)

